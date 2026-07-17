# -*- coding: utf-8 -*-
"""
AIOps 智能运维平台 PPT 重建脚本
使用 python-pptx 重建 30 页可编辑 PPT
"""
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ============================================================
# 颜色常量
# ============================================================
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG  = RGBColor(0xF5, 0xF7, 0xFA)
C_BORDER   = RGBColor(0xE0, 0xE0, 0xE0)
C_BLUE     = RGBColor(0x15, 0x65, 0xC0)
C_GREEN    = RGBColor(0x00, 0x89, 0x7B)
C_ORANGE   = RGBColor(0xFF, 0x6F, 0x00)
C_RED      = RGBColor(0xE5, 0x39, 0x35)
C_SUCCESS  = RGBColor(0x43, 0xA0, 0x47)
C_DARK     = RGBColor(0x21, 0x21, 0x21)
C_GRAY     = RGBColor(0x61, 0x61, 0x61)
C_LGRAY    = RGBColor(0x9E, 0x9E, 0x9E)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_COVER_BG = RGBColor(0x0D, 0x47, 0xA1)

# 浅色背景（用于 fill-opacity 模拟）
C_BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFB)
C_GREEN_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)
C_ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE0)
C_RED_LIGHT = RGBColor(0xFF, 0xEB, 0xEE)

# ============================================================
# 辅助函数
# ============================================================
def add_text(slide, text, left_pt, top_pt, width_pt, height_pt,
             size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    """添加文本框"""
    tb = slide.shapes.add_textbox(
        Inches(left_pt/960), Inches(top_pt/960),
        Inches(width_pt/960), Inches(height_pt/960))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    return tb

def add_rect(slide, left_pt, top_pt, width_pt, height_pt,
             fill=None, line=None, line_width=1, rx=0):
    """添加矩形（圆角用 rounded_rectangle）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE,
        Inches(left_pt/960), Inches(top_pt/960),
        Inches(width_pt/960), Inches(height_pt/960))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_footer(slide, page_num, total=30):
    add_text(slide, f"{page_num} / {total}", 60, 700, 200, 20, size=11, color=C_LGRAY)
    add_text(slide, "AIOps 智能运维平台", 1100, 700, 180, 20,
             size=11, color=C_LGRAY, align=PP_ALIGN.RIGHT)

def add_header(slide, title, color=C_BLUE):
    add_text(slide, title, 60, 95, 1160, 50, size=32, bold=True, color=color)
    add_rect(slide, 60, 108, 80, 4, fill=color, rx=2)

def add_note_bar(slide, text, color=C_BLUE):
    """底部注释条"""
    add_rect(slide, 60, 450, 1160, 50, fill=C_BLUE_LIGHT, rx=8)
    add_text(slide, text, 640, 465, 600, 30, size=16, color=color, align=PP_ALIGN.CENTER)

def add_arrow_triangle(slide, left_pt, top_pt, color, rotation=0):
    """添加三角形箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(left_pt/960), Inches(top_pt/960),
        Inches(20/960), Inches(15/960))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = rotation
    return shape

def add_circle_with_text(slide, cx, cy, r, fill_color, text, text_color=C_WHITE, text_size=20):
    """添加圆形+文字"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches((cx-r)/960), Inches((cy-r)/960),
        Inches((r*2)/960), Inches((r*2)/960))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(text_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "微软雅黑"
    return shape

def add_card(slide, left, top, width, height, fill=C_CARD_BG, line=C_BORDER, rx=12):
    """通用卡片"""
    return add_rect(slide, left, top, width, height, fill=fill, line=line, line_width=1, rx=rx)

def add_flow_box(slide, left, top, width, height, title, subtitle=None,
                 fill=C_CARD_BG, line_color=C_BLUE, rx=12):
    """流程方框"""
    shape = add_rect(slide, left, top, width, height, fill=fill, line=line_color, line_width=1.5, rx=rx)
    # 标题
    title_size = 18 if len(title) <= 8 else 16
    add_text(slide, title, left+10, top+15, width-20, 35, size=title_size, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    if subtitle:
        sub_size = 14 if len(subtitle) <= 15 else 12
        add_text(slide, subtitle, left+10, top+50, width-20, 30, size=sub_size, color=C_GRAY, align=PP_ALIGN.CENTER)
    return shape

# ============================================================
# 创建 PPT
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# Slide 01 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
# 深蓝背景
add_rect(slide, 0, 0, 1280, 720, fill=C_COVER_BG)
# 主标题
add_text(slide, "AIOps 智能运维平台", 0, 240, 1280, 80, size=54, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)
# 副标题
add_text(slide, "系统架构数据流交互图详解", 0, 346, 1280, 50, size=28,
         color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)
# 信息
add_text(slide, "2026-07-15  |  客户技术交流", 0, 422, 1280, 30, size=16,
         color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
add_footer(slide, 1)

# ============================================================
# Slide 02 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_text(slide, "目录", 60, 95, 200, 50, size=32, bold=True, color=C_BLUE)
add_rect(slide, 60, 108, 80, 4, fill=C_BLUE, rx=2)

# 四卡片 2x2
cards_data = [
    (60, 150, "核心观测闭环", "#1565C0", "资产发现 → 指标采集 → 异常检测\n→ 告警生成 → 移动端推送", "链路 1 - 4"),
    (640, 150, "智能诊断与知识", "#00897B", "变更关联 → RCA 根因 → 知识沉淀\n→ AI 问答 → SOP 生成 → 审批流", "链路 5 - 12"),
    (60, 400, "自动化与自愈", "#FF6F00", "告警触发巡检 → 混沌稳态验证\n→ 自动回滚 → On-Call 轮转", "链路 13 - 15"),
    (640, 400, "AI 增强与运营", "#E53935", "Agent 评估 → 异常回测 → 运营飞轮\n→ 诊断工具 → 仪表盘 → 移动端", "链路 16 - 24"),
]

for left, top, title, color_hex, content, tag in cards_data:
    fill_color = RGBColor(int(color_hex[1:3],16), int(color_hex[3:5],16), int(color_hex[5:7],16))
    card = add_card(slide, left, top, 540, 220)
    # 圆形编号
    circle_x = left + 50
    circle_y = top + 50
    add_circle_with_text(slide, circle_x, circle_y, 24, fill_color, title[0])
    # 标题
    add_text(slide, title, left+80, top+35, 400, 35, size=22, bold=True, color=fill_color)
    # 分割线
    add_rect(slide, left+24, top+90, 492, 1, fill=C_BORDER)
    # 内容
    lines = content.split('\n')
    for i, line in enumerate(lines):
        add_text(slide, line, left+40, top+110+i*30, 460, 25, size=18, color=C_GRAY)
    # 标签
    tag_color = fill_color
    add_rect(slide, left+40, top+185, 180, 30, fill=RGBColor(
        int(color_hex[1:3],16), int(color_hex[3:5],16), int(color_hex[5:7],16)) & RGBColor(0xFF,0xFF,0xFF) if False else fill_color,
             rx=6)
    # 用浅色背景
    light_color = C_BLUE_LIGHT if fill_color == C_BLUE else \
                  C_GREEN_LIGHT if fill_color == C_GREEN else \
                  C_ORANGE_LIGHT if fill_color == C_ORANGE else C_RED_LIGHT
    add_rect(slide, left+40, top+185, 180, 30, fill=light_color, rx=6)
    add_text(slide, tag, left+40, top+188, 180, 26, size=15, bold=True, color=fill_color, align=PP_ALIGN.CENTER)

add_footer(slide, 2)

# ============================================================
# Slide 03 平台架构总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "平台架构总览", C_BLUE)

layers = [
    ("资产层 | CMDB / 拓扑发现", RGBColor(0x9E, 0x9E, 0x9E)),
    ("采集层 | SSH / K8s API / Metricbeat", RGBColor(0x9E, 0x9E, 0x9E)),
    ("存储层 | VictoriaMetrics 双写 / SQLite", C_BLUE),
    ("检测层 | 7 种异常算法 / 统计推断", C_BLUE),
    ("告警层 | 聚合防抖 / 升级 / 静默", C_ORANGE),
    ("诊断层 | 6 种 RCA / 知识推理", C_GREEN),
    ("知识层 | RAG 检索 / SOP 生成", C_GREEN),
    ("自愈层 | 规则引擎 / 混沌验证", C_RED),
]

y = 150
for i, (text, color) in enumerate(layers):
    add_rect(slide, 60, y, 1160, 80, fill=color, rx=4)
    add_text(slide, text, 60, y+20, 1160, 50, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    y += 92

add_footer(slide, 3)

# ============================================================
# Slide 04 核心数据流（上）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "核心数据流（上）：观测 → 检测 → 告警", C_BLUE)

nodes = [
    ("资产自动发现", "TCP/SSH/ICMP", C_BLUE),
    ("指标采集", "SSH 19 命令", C_BLUE),
    ("VM 双写", "VictoriaMetrics", C_GREEN),
    ("健康引擎", "三层判定", C_GREEN),
    ("可观测展示", "灭火图/仪表盘", C_BLUE),
    ("异常检测", "7 种算法", C_ORANGE),
]

x = 85
for i, (title, sub, color) in enumerate(nodes):
    add_flow_box(slide, x, 280, 160, 100, title, sub, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+165, 318, color)
    x += 190

add_note_bar(slide, "链路 1 前半段  |  资产自动发现 → 指标采集 → 时序存储 → 健康判定 → 可观测展示 → 异常检测")
add_text(slide, "从资产自动发现到指标采集，通过 VictoriaMetrics 双写实现时序数据存储，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "健康引擎进行三层判定，结果可视化展示，异常检测算法识别故障。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 4)

# ============================================================
# Slide 05 核心数据流（下）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "核心数据流（下）：诊断 → 知识 → AI → 自愈 → SRE", C_BLUE)

nodes = [
    ("告警生成", "聚合防抖", C_ORANGE),
    ("通知推送", "移动端/邮件", C_BLUE),
    ("故障关联RCA", "6种根因", C_RED),
    ("知识沉淀", "RAG索引", C_GREEN),
    ("AI助手", "LLM对话", C_BLUE),
    ("自愈执行", "规则引擎", C_ORANGE),
    ("SRE闭环", "复盘改进", C_GREEN),
]

x = 80
for i, (title, sub, color) in enumerate(nodes):
    add_flow_box(slide, x, 280, 150, 100, title, sub, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+155, 318, color)
    x += 178

add_note_bar(slide, "链路 1 后半段  |  告警生成 → 通知推送 → 故障关联RCA → 知识沉淀 → AI助手 → 自愈执行 → SRE闭环")
add_text(slide, "告警生成后经过多渠道通知，RCA根因分析驱动知识沉淀，AI助手提供决策支持，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "自愈执行自动恢复故障，SRE闭环实现持续改进。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 5)

# ============================================================
# Slide 06 K8s 指标与事件接入
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "K8s 指标与事件接入", C_BLUE)

# 左侧卡片
add_card(slide, 60, 150, 540, 420)
add_text(slide, "链路 2  K8s指标采集", 84, 172, 492, 30, size=18, bold=True, color=C_BLUE)
items_l = [
    ("K8s Metrics API 轮询", C_BLUE),
    ("采集 CPU / MEM / 网络", C_BLUE),
    ("生成 MetricRecord", C_BLUE),
    ("双写 VictoriaMetrics", C_GREEN),
]
y = 220
for title, color in items_l:
    add_flow_box(slide, 100, y, 400, 70, title, line_color=color)
    if y < 400:
        add_arrow_triangle(slide, 300, y+75, C_GREEN, rotation=90)
    y += 95

# 右侧卡片
add_card(slide, 640, 150, 580, 420)
add_text(slide, "链路 3  K8s事件告警", 664, 172, 532, 30, size=18, bold=True, color=C_ORANGE)
items_r = [
    ("K8s Events 轮询", C_ORANGE),
    ("Warning / Error 过滤", C_ORANGE),
    ("Alert 映射规则", C_ORANGE),
    ("防抖去重聚合", C_GREEN),
]
y = 220
for title, color in items_r:
    add_flow_box(slide, 680, y, 460, 70, title, line_color=color)
    if y < 400:
        add_arrow_triangle(slide, 880, y+75, C_GREEN, rotation=90)
    y += 95

add_footer(slide, 6)

# ============================================================
# Slide 07 移动端告警实时推送
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "移动端告警实时推送", C_BLUE)

nodes = [
    ("WebSocket连接", C_BLUE),
    ("Token鉴权", C_BLUE),
    ("实时推送", C_GREEN),
    ("分栏展示", C_GREEN),
    ("操作回写", C_ORANGE),
]

x = 100
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 245

add_note_bar(slide, "链路 4  |  WebSocket连接 → Token鉴权 → 实时推送 → 分栏展示 → 操作回写")
add_text(slide, "通过 WebSocket 保持长连接，Token 鉴权确保安全，告警实时推送至移动端，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "分栏展示便于查看，操作结果自动回写至平台。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 7)

# ============================================================
# Slide 08 事件变更关联故障
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "事件变更关联故障", C_BLUE)

nodes = [
    ("变更审批单", C_BLUE),
    ("30分钟窗口", C_ORANGE),
    ("关联Incident", C_RED),
    ("RCA纳入", C_GREEN),
]

x = 100
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 250, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+255, 302, color)
    x += 300

add_note_bar(slide, "链路 5  |  变更审批单 → 30分钟窗口 → 关联Incident → RCA纳入")
add_text(slide, "变更审批单触发后，在30分钟窗口期内自动关联相关故障事件，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "将变更作为根因分析的重要因素纳入 RCA 分析。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 8)

# ============================================================
# Slide 09 知识指导SRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "知识指导SRE", C_BLUE)

# 左侧
add_card(slide, 60, 150, 540, 420)
add_text(slide, "知识库 → SRE", 84, 172, 492, 30, size=18, bold=True, color=C_GREEN)
items_l = [
    ("RAG检索知识", C_GREEN),
    ("故障案例匹配", C_GREEN),
    ("SOP推荐", C_GREEN),
    ("执行指引", C_BLUE),
]
y = 220
for title, color in items_l:
    add_flow_box(slide, 100, y, 400, 70, title, line_color=color)
    if y < 400:
        add_arrow_triangle(slide, 300, y+75, C_GREEN, rotation=90)
    y += 95

# 右侧
add_card(slide, 640, 150, 580, 420)
add_text(slide, "SRE面板", 664, 172, 532, 30, size=18, bold=True, color=C_BLUE)
items_r = [
    ("值班人员", C_BLUE),
    ("告警面板", C_BLUE),
    ("响应记录", C_ORANGE),
    ("复盘报告", C_GREEN),
]
y = 220
for title, color in items_r:
    add_flow_box(slide, 680, y, 460, 70, title, line_color=color)
    if y < 400:
        add_arrow_triangle(slide, 880, y+75, C_BLUE, rotation=90)
    y += 95

add_footer(slide, 9)

# ============================================================
# Slide 10 SLO违规触发自愈
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "SLO违规触发自愈", C_BLUE)

nodes = [
    ("SLO计算", C_BLUE),
    ("违规告警", C_ORANGE),
    ("规则匹配", C_RED),
    ("扩缩容重启", C_ORANGE),
    ("30分钟追踪", C_GREEN),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 6  |  SLO计算 → 违规告警 → 规则匹配 → 扩缩容重启 → 30分钟追踪")
add_text(slide, "实时计算 SLO 可用性，当违规触发告警时，自动匹配自愈规则执行扩缩容或重启，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "追踪30分钟内的恢复效果，记录自愈执行日志。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 10)

# ============================================================
# Slide 11 告警直接触发自动响应
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "告警直接触发自动响应", C_BLUE)

nodes = [
    ("定时扫描", C_BLUE),
    ("SSH执行", C_ORANGE),
    ("黑名单校验", C_RED),
    ("RemediationLog", C_GREEN),
    ("30分钟追踪", C_GREEN),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 7  |  定时扫描 → SSH执行 → 黑名单校验 → RemediationLog → 30分钟追踪")
add_text(slide, "定时扫描资产指标，当告警触发时，通过 SSH 执行预设修复命令，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "黑名单校验防止危险操作，全程记录 RemediationLog 并追踪30分钟恢复效果。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 11)

# ============================================================
# Slide 12 故障自动沉淀知识
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "故障自动沉淀知识", C_BLUE)

nodes = [
    ("故障resolved", C_GREEN),
    ("LLM生成", C_BLUE),
    ("RAG索引", C_GREEN),
    ("TF-IDF", C_ORANGE),
    ("检索召回", C_BLUE),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 8  |  故障resolved → LLM生成 → RAG索引 → TF-IDF → 检索召回")
add_text(slide, "故障解决后，LLM 自动生成故障分析报告，通过 RAG 索引和 TF-IDF 关键词提取，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "将知识沉淀至知识库，支持后续相似故障的快速检索和召回。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 12)

# ============================================================
# Slide 13 AI检索知识
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "AI检索知识", C_BLUE)

nodes = [
    ("用户提问", C_BLUE),
    ("LLM tool_call", C_BLUE),
    ("RAG检索", C_GREEN),
    ("结果注入", C_GREEN),
    ("采纳草稿", C_ORANGE),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 9  |  用户提问 → LLM tool_call → RAG检索 → 结果注入 → 采纳草稿")
add_text(slide, "用户提问时，LLM 通过 tool_call 调用 RAG 检索相关知识，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "将检索结果注入上下文，生成回答供用户采纳并形成知识草稿。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 13)

# ============================================================
# Slide 14 故障审批流状态机
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "故障审批流状态机", C_BLUE)

# 三个状态
states = [
    ("open", "故障已创建", C_RED),
    ("pending_approval", "审批中", C_ORANGE),
    ("resolved", "已解决", C_SUCCESS),
]

x = 120
for i, (state, desc, color) in enumerate(states):
    add_rect(slide, x, 250, 300, 150, fill=color, rx=12)
    add_text(slide, state, x, 280, 300, 50, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x, 340, 300, 40, size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        add_arrow_triangle(slide, x+310, 312, C_DARK)
    x += 380

add_note_bar(slide, "链路 10  |  故障审批流状态机：open → pending_approval → resolved")
add_text(slide, "故障单创建后状态为 open，提交审批后转为 pending_approval，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "审批通过或驳回后流转至 resolved 状态，形成完整的审批闭环。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 14)

# ============================================================
# Slide 15 告警静默管理
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "告警静默管理", C_BLUE)

nodes = [
    ("创建静默规则", C_BLUE),
    ("批量silenced", C_ORANGE),
    ("过期恢复", C_GREEN),
    ("手动取消", C_RED),
]

x = 100
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 250, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+255, 302, color)
    x += 300

add_note_bar(slide, "链路 11  |  创建静默规则 → 批量silenced → 过期恢复 → 手动取消")
add_text(slide, "支持创建静默规则对告警进行批量静默，过期后自动恢复，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "也可手动取消静默，实现灵活的告警管理。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 15)

# ============================================================
# Slide 16 告警触发巡检
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "告警触发巡检", C_BLUE)

nodes = [
    ("告警触发", C_RED),
    ("匹配模板", C_BLUE),
    ("8类检查", C_ORANGE),
    ("AI报告", C_GREEN),
    ("关联告警", C_BLUE),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 13  |  告警触发 → 匹配模板 → 8类检查 → AI报告 → 关联告警")
add_text(slide, "告警触发后自动匹配巡检模板，执行8类检查项（CPU/内存/磁盘/告警/状态等），", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "AI 生成巡检报告，并将结果关联至原始告警。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 16)

# ============================================================
# Slide 17 混沌工程
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "混沌工程", C_BLUE)

nodes = [
    ("采集基线", C_BLUE),
    ("注入故障", C_RED),
    ("sleep观察", C_ORANGE),
    ("采集对比", C_BLUE),
    ("稳态验证", C_GREEN),
    ("自动回滚", C_RED),
    ("结果标记", C_GREEN),
]

x = 75
colors = [C_BLUE, C_RED, C_ORANGE, C_BLUE, C_GREEN, C_RED, C_GREEN]
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 140, 115, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+145, 305, color)
    x += 167

add_note_bar(slide, "链路 14  |  采集基线 → 注入故障 → sleep观察 → 采集对比 → 稳态验证 → 自动回滚 → 结果标记")
add_text(slide, "混沌工程通过注入故障验证系统韧性，采集稳态基线数据，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "观察系统表现，验证稳态假设，自动回滚并标记实验结果。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 17)

# ============================================================
# Slide 18 On-Call自动轮转
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "On-Call自动轮转", C_BLUE)

nodes = [
    ("is_auto_rotate", C_BLUE),
    ("周期过期", C_ORANGE),
    ("自动轮转", C_GREEN),
    ("手动触发", C_ORANGE),
    ("handover", C_BLUE),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 15  |  is_auto_rotate → 周期过期 → 自动轮转 → 手动触发 → handover")
add_text(slide, "On-Call 排班支持自动轮转，周期过期后自动切换值班人员，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "也支持手动触发交接班，确保告警响应无缝衔接。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 18)

# ============================================================
# Slide 19 Agent评估与AB测试
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "Agent评估与AB测试", C_BLUE)

# 左侧
add_card(slide, 60, 150, 540, 420)
add_text(slide, "AgentEvaluation 评估指标", 84, 172, 492, 30, size=18, bold=True, color=C_BLUE)
items_l = [
    ("对话轮次", C_BLUE),
    ("响应时长", C_BLUE),
    ("解决率", C_GREEN),
    ("满意度", C_GREEN),
]
y = 220
for title, color in items_l:
    add_flow_box(slide, 100, y, 400, 70, title, line_color=color)
    if y < 400:
        add_arrow_triangle(slide, 300, y+75, C_BLUE, rotation=90)
    y += 95

# 右侧
add_card(slide, 640, 150, 580, 420)
add_text(slide, "A/B测试流程", 664, 172, 532, 30, size=18, bold=True, color=C_ORANGE)
items_r = [
    ("流量分配", C_ORANGE),
    ("策略A", C_BLUE),
    ("策略B", C_GREEN),
    ("效果对比", C_RED),
]
y = 220
for title, color in items_r:
    add_flow_box(slide, 680, y, 460, 70, title, line_color=color)
    if y < 400:
        add_arrow_triangle(slide, 880, y+75, C_ORANGE, rotation=90)
    y += 95

add_footer(slide, 19)

# ============================================================
# Slide 20 异常检测回测
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "异常检测回测", C_BLUE)

nodes = [
    ("标注数据", C_BLUE),
    ("F1评估", C_GREEN),
    ("算法推荐", C_ORANGE),
    ("特征匹配回退", C_RED),
    ("统计推断", C_BLUE),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 17  |  标注数据 → F1评估 → 算法推荐 → 特征匹配回退 → 统计推断")
add_text(slide, "通过标注数据对异常检测算法进行 F1 评估，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "推荐最优算法，特征匹配失败时回退至统计推断方法。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 20)

# ============================================================
# Slide 21 SOP知识生成
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "SOP知识生成", C_BLUE)

nodes = [
    ("故障详情", C_BLUE),
    ("LLM生成SOP", C_GREEN),
    ("草稿审批", C_ORANGE),
    ("正式知识库", C_BLUE),
]

x = 100
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 250, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+255, 302, color)
    x += 300

add_note_bar(slide, "链路 18  |  故障详情 → LLM生成SOP → 草稿审批 → 正式知识库")
add_text(slide, "故障解决后，LLM 自动生成 SOP 操作步骤，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "经过草稿审批流程，最终沉淀至正式知识库。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 21)

# ============================================================
# Slide 22 运营数据飞轮
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "运营数据飞轮", C_BLUE)

# 6卡片 2x3
cards_22 = [
    ("MTTA 平均响应时间", "3.2", "分钟", "↓ 42%", C_BLUE),
    ("MTTR 平均修复时间", "12.5", "分钟", "↓ 38%", C_GREEN),
    ("自愈成功率", "68%", "", "↑ 15%", C_ORANGE),
    ("AI 对话效能", "1,247", "次/周", "↑ 28%", C_BLUE),
    ("通知推送统计", "8,532", "条/月", "", C_GREEN),
    ("工具使用排行", "Top3", "", "", C_ORANGE),
]

positions = [(60,140), (460,140), (860,140), (60,325), (460,325), (860,325)]

for idx, (label, num, unit, trend, color) in enumerate(cards_22):
    x, y = positions[idx]
    add_card(slide, x, y, 360, 155)
    add_text(slide, label, x+24, y+32, 312, 25, size=16, color=C_GRAY)
    add_text(slide, num, x+24, y+90, 120, 50, size=48, bold=True, color=color)
    if unit:
        add_text(slide, unit, x+140, y+100, 80, 30, size=20, color=C_GRAY)
    if trend:
        # 绿色三角箭头（向上）
        if "↑" in trend:
            add_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches((x+260)/960), Inches((y+95)/960),
                Inches(15/960), Inches(12/960))
            add_shape.fill.solid()
            add_shape.fill.fore_color.rgb = C_SUCCESS
            add_shape.line.fill.background()
            add_shape.rotation = 270
        add_text(slide, trend, x+240, y+115, 100, 25, size=14, color=C_SUCCESS)

# 底部注释条
add_rect(slide, 60, 510, 1160, 50, fill=C_BLUE_LIGHT, rx=8)
add_text(slide, "链路 19  |  运营数据飞轮 — 6 大 KPI 持续度量平台效能",
         640, 525, 600, 30, size=16, color=C_BLUE, align=PP_ALIGN.CENTER)

add_text(slide, "运营数据飞轮持续度量 MTTA/MTTR/自愈成功率等关键指标，", 60, 590, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "数据驱动平台迭代方向，实现效能持续提升。", 60, 615, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 22)

# ============================================================
# Slide 23 分层诊断工具
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "分层诊断工具", C_BLUE)

# 三层横向大条
tools_23 = [
    ("Snapshot 快照层（6工具）", "进程上下文 | 端口扫描 | 日志 tail | 网络连接 | 磁盘IO | 内存 dump",
     RGBColor(0x15, 0x65, 0xC0)),
    ("Focused 定向层（12工具）", "CPU Profiler | 网络抓包 | 链路追踪 | MySQL诊断 | Redis诊断 | JVM分析...",
     RGBColor(0x00, 0x89, 0x7B)),
    ("Flexible 灵活层（2工具）", "自由脚本执行 | 远程终端",
     RGBColor(0xFF, 0x6F, 0x00)),
]

y = 160
for text, tools, color in tools_23:
    add_rect(slide, 60, y, 1160, 160, fill=color, rx=8)
    add_text(slide, text, 84, y+20, 1092, 40, size=24, bold=True, color=C_WHITE)
    add_text(slide, tools, 84, y+75, 1092, 60, size=16, color=C_WHITE)
    y += 185

add_footer(slide, 23)

# ============================================================
# Slide 24 仪表盘拖拽
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "仪表盘拖拽", C_BLUE)

# 左侧 16种卡片库
add_card(slide, 60, 150, 540, 480)
add_text(slide, "16种卡片库", 84, 172, 492, 30, size=18, bold=True, color=C_BLUE)

card_types = [
    "折线图", "柱状图", "饼图", "仪表盘",
    "数字卡", "趋势卡", "排行榜", "状态卡",
    "地图", "热力图", "雷达图", "漏斗图",
    "表格", "日志流", "告警列表", "链路图"
]
for i, ct in enumerate(card_types):
    row = i // 4
    col = i % 4
    x = 80 + col * 125
    y = 220 + row * 100
    add_rect(slide, x, y, 115, 80, fill=C_CARD_BG, line=C_BORDER, rx=6)
    add_text(slide, ct, x, y+25, 115, 30, size=14, color=C_DARK, align=PP_ALIGN.CENTER)

# 右侧
add_card(slide, 640, 150, 580, 480)
add_text(slide, "gridstack 画布功能", 664, 172, 532, 30, size=18, bold=True, color=C_GREEN)
features = [
    "拖拽调整位置和大小",
    "网格化布局系统",
    "实时预览效果",
    "保存自定义布局",
    "多套预设模板",
    "响应式自适应"
]
y = 220
for f in features:
    add_rect(slide, 680, y, 500, 55, fill=C_GREEN_LIGHT, rx=6)
    add_text(slide, f, 700, y+12, 460, 35, size=16, color=C_GREEN)
    y += 72

add_footer(slide, 24)

# ============================================================
# Slide 25 资产健康度
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "资产健康度", C_BLUE)

nodes = [
    ("CPU/MEM/DISK+告警", C_RED),
    ("0-100评分", C_BLUE),
    ("状态变化", C_ORANGE),
    ("AssetChangeLog", C_GREEN),
    ("分阈值告警", C_ORANGE),
]

x = 80
for i, (title, color) in enumerate(nodes):
    add_flow_box(slide, x, 260, 200, 110, title, line_color=color)
    if i < len(nodes)-1:
        add_arrow_triangle(slide, x+205, 302, color)
    x += 240

add_note_bar(slide, "链路 22  |  CPU/MEM/DISK+告警 → 0-100评分 → 状态变化 → AssetChangeLog → 分阈值告警")
add_text(slide, "综合 CPU、内存、磁盘指标及活跃告警，对资产进行 0-100 健康度评分，", 60, 520, 1160, 25, size=16, color=C_GRAY)
add_text(slide, "状态变化记录至变更日志，触发分阈值告警。", 60, 550, 1160, 25, size=16, color=C_GRAY)
add_footer(slide, 25)

# ============================================================
# Slide 26 移动端运维闭环
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "移动端运维闭环", C_BLUE)

cards_26 = [
    ("日志搜索", "链路 23-1", C_BLUE),
    ("告警批量操作", "链路 23-2", C_ORANGE),
    ("故障单创建", "链路 23-3", C_RED),
    ("AI会话列表", "链路 23-4", C_GREEN),
]

for idx, (title, tag, color) in enumerate(cards_26):
    row = idx // 2
    col = idx % 2
    x = 60 + col * 610
    y = 160 + row * 250
    add_card(slide, x, y, 560, 220)
    add_rect(slide, x+24, y+24, 100, 35, fill=color, rx=6)
    add_text(slide, tag, x+24, y+28, 100, 28, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x+24, y+80, 512, 40, size=24, bold=True, color=color)

add_footer(slide, 26)

# ============================================================
# Slide 27 Agent能力中心
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "Agent能力中心", C_BLUE)

# 左侧工具列表
add_card(slide, 60, 150, 540, 480)
add_text(slide, "工具列表（41个）", 84, 172, 492, 30, size=18, bold=True, color=C_BLUE)
tools = [
    "query_logs | 日志查询 | read_only",
    "query_traces | 链路追踪 | read_only",
    "execute_command | 命令执行 | critical",
    "query_mysql | 数据库查询 | high",
    "restart_service | 服务重启 | critical",
]
y = 220
for t in tools:
    parts = t.split(" | ")
    add_rect(slide, 80, y, 480, 50, fill=C_CARD_BG, rx=4)
    add_text(slide, parts[0], 90, y+10, 150, 30, size=14, bold=True, color=C_DARK)
    add_text(slide, parts[1], 240, y+10, 150, 30, size=14, color=C_GRAY)
    add_text(slide, parts[2], 400, y+10, 140, 30, size=14, color=C_ORANGE)
    y += 60

# 右侧统计面板
add_card(slide, 640, 150, 580, 480)
add_text(slide, "统计面板", 664, 172, 532, 30, size=18, bold=True, color=C_GREEN)
stats = [
    ("read_only 工具", "20", C_BLUE),
    ("low 风险", "4", C_GREEN),
    ("medium 风险", "6", C_ORANGE),
    ("high 风险", "5", C_RED),
    ("critical 风险", "3", RGBColor(0xB7, 0x00, 0x00)),
    ("advisory", "3", C_GRAY),
]
y = 220
for name, count, color in stats:
    add_rect(slide, 680, y, 480, 55, fill=C_CARD_BG, rx=4)
    add_text(slide, name, 700, y+12, 300, 30, size=16, color=C_DARK)
    add_text(slide, count, 1050, y+5, 80, 40, size=28, bold=True, color=color)
    y += 68

add_footer(slide, 27)

# ============================================================
# Slide 28 平台核心价值
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_text(slide, "平台核心价值", 60, 95, 200, 50, size=32, bold=True, color=C_BLUE)
add_rect(slide, 60, 108, 80, 4, fill=C_BLUE, rx=2)

cards_28 = [
    ("全观测闭环", "#1565C0", "资产→采集→检测→告警→诊断→自愈", "链路 1 全流程"),
    ("AI 智能诊断", "#00897B", "7 种异常算法 + RAG + 6 种 RCA", "链路 9-14"),
    ("自动化自愈", "#FF6F00", "规则引擎 + 效果追踪 + 混沌回滚", "链路 7-8"),
    ("知识飞轮", "#E53935", "自动沉淀 + 审批 + 检索 + 重排", "链路 9-13"),
]

for idx, (title, color_hex, desc, tag) in enumerate(cards_28):
    row = idx // 2
    col = idx % 2
    x = 60 + col * 610
    y = 160 + row * 250
    fill_color = RGBColor(int(color_hex[1:3],16), int(color_hex[3:5],16), int(color_hex[5:7],16))
    add_card(slide, x, y, 540, 220)
    add_text(slide, title, x+24, y+24, 492, 35, size=22, bold=True, color=fill_color)
    add_rect(slide, x+24, y+75, 492, 1, fill=C_BORDER)
    add_text(slide, desc, x+40, y+95, 460, 50, size=18, color=C_GRAY)
    light_color = C_BLUE_LIGHT if fill_color == C_BLUE else \
                  C_GREEN_LIGHT if fill_color == C_GREEN else \
                  C_ORANGE_LIGHT if fill_color == C_ORANGE else C_RED_LIGHT
    add_rect(slide, x+40, y+165, 180, 30, fill=light_color, rx=6)
    add_text(slide, tag, x+40, y+168, 180, 26, size=15, bold=True, color=fill_color, align=PP_ALIGN.CENTER)

add_footer(slide, 28)

# ============================================================
# Slide 29 技术亮点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_BG)
add_header(slide, "技术亮点", C_BLUE)

highlights = [
    ("VictoriaMetrics 时序存储", "高性能时序数据库，支持高并发写入和压缩存储，\n双写机制保障数据可靠性。", RGBColor(0x15, 0x65, 0xC0)),
    ("BGE-M3 RAG 检索增强", "新一代向量检索技术，多向量融合，\n支持密集检索、稀疏检索和重排序。", RGBColor(0x00, 0x89, 0x7B)),
    ("混沌稳态验证", "通过故障注入验证系统韧性，\n稳态假设验证保障自愈可靠性。", RGBColor(0xFF, 0x6F, 0x00)),
]

y = 160
for title, desc, color in highlights:
    add_rect(slide, 60, y, 1160, 160, fill=color, rx=8)
    add_text(slide, title, 84, y+20, 1092, 45, size=26, bold=True, color=C_WHITE)
    add_text(slide, desc, 84, y+75, 1092, 60, size=16, color=C_WHITE)
    y += 185

add_footer(slide, 29)

# ============================================================
# Slide 30 结尾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 1280, 720, fill=C_COVER_BG)
add_text(slide, "感谢聆听", 0, 300, 1280, 80, size=54, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)
add_footer(slide, 30)

# ============================================================
# 保存 PPT
# ============================================================
output_path = r"D:\AIOPS\project08\docs\AIOps系统架构交互图_客户交流PPT_20260715_可编辑版.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
print(f"共生成 {len(prs.slides)} 页")
